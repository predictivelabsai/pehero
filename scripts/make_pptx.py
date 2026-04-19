"""Generate a PEHero product-tour .pptx — editable slide deck for sharing.

Mirrors the content of scripts/make_pdf.py but writes to docs/pehero-product-tour.pptx
so the deck can be edited in PowerPoint / Keynote / Google Slides.

Usage:
    python -m scripts.make_pptx
"""

from __future__ import annotations

from pathlib import Path

from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Cm, Inches, Pt


ROOT = Path(__file__).resolve().parents[1]
SHOTS = ROOT / "screenshots"
OUT = ROOT / "docs" / "pehero-product-tour.pptx"

# 16:9 slide — 33.87 × 19.05 cm
SLIDE_W = Cm(33.87)
SLIDE_H = Cm(19.05)

# PEHero palette (parchment + forest)
BG         = RGBColor(0xF7, 0xF6, 0xF1)
INK        = RGBColor(0x14, 0x23, 0x1B)
INK_MUTED  = RGBColor(0x41, 0x50, 0x46)
INK_DIM    = RGBColor(0x7A, 0x86, 0x7E)
ACCENT     = RGBColor(0x1F, 0x5D, 0x43)
ACCENT_DIM = RGBColor(0xCF, 0xE5, 0xDA)
RULE       = RGBColor(0xE3, 0xDF, 0xD2)


def _blank(prs: Presentation):
    """Add a blank slide with the parchment background + footer rule."""
    layout = prs.slide_layouts[6]  # blank
    slide = prs.slides.add_slide(layout)
    bg = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H
    )
    bg.line.fill.background()
    bg.fill.solid()
    bg.fill.fore_color.rgb = BG
    bg.shadow.inherit = False

    # footer line
    footer_rule = slide.shapes.add_connector(
        1,  # straight line
        Cm(1.5), SLIDE_H - Cm(1.5),
        SLIDE_W - Cm(1.5), SLIDE_H - Cm(1.5),
    )
    footer_rule.line.color.rgb = RULE
    footer_rule.line.width = Pt(0.5)

    # footer text
    tb = slide.shapes.add_textbox(Cm(1.5), SLIDE_H - Cm(1.2),
                                  SLIDE_W - Cm(3), Cm(0.6))
    tf = tb.text_frame
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    r = p.add_run()
    r.text = "PEHero · Your Private Equity AI Agent Squad"
    r.font.size = Pt(8)
    r.font.name = "Inter"
    r.font.color.rgb = INK_DIM

    return slide


def _text(slide, x, y, w, h, text, *, size=11, bold=False, color=INK,
          align=PP_ALIGN.LEFT, font="Inter"):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.name = font
    r.font.color.rgb = color
    return tb


def _bullets(slide, x, y, w, h, items: list[str], *, size=12):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.space_after = Pt(6)
        r = p.add_run()
        r.text = f"•  {item}"
        r.font.size = Pt(size)
        r.font.name = "Inter"
        r.font.color.rgb = INK
    return tb


def _fit_image_inches(path: Path, max_w_in: float, max_h_in: float) -> tuple[float, float]:
    img = Image.open(path)
    w, h = img.size
    # 96 dpi assumption
    w_in = w / 96.0
    h_in = h / 96.0
    scale = min(max_w_in / w_in, max_h_in / h_in)
    return w_in * scale, h_in * scale


def _hero_slide(prs):
    s = _blank(prs)
    # Large centred title
    _text(s, Cm(0), Cm(6.5), SLIDE_W, Cm(3.5),
          "PEHero", size=54, bold=True, color=INK, align=PP_ALIGN.CENTER)
    _text(s, Cm(0), Cm(10.5), SLIDE_W, Cm(1.5),
          "Your Private Equity AI Agent Squad.",
          size=20, color=INK_MUTED, align=PP_ALIGN.CENTER)
    _text(s, Cm(0), Cm(12.2), SLIDE_W, Cm(1),
          "Sourcing · underwriting · diligence · capital · portfolio operations",
          size=14, color=INK_DIM, align=PP_ALIGN.CENTER)
    return s


def _agenda_slide(prs):
    s = _blank(prs)
    _text(s, Cm(1.5), Cm(1.5), SLIDE_W - Cm(3), Cm(0.8),
          "TOUR", size=10, bold=True, color=ACCENT, align=PP_ALIGN.LEFT)
    _text(s, Cm(1.5), Cm(2.2), SLIDE_W - Cm(3), Cm(1.6),
          "What you'll see", size=28, bold=True, color=INK)
    _text(s, Cm(1.5), Cm(3.9), SLIDE_W - Cm(3), Cm(1),
          "Six product surfaces, in order of a deal's life.",
          size=14, color=INK_MUTED)

    rows = [
        ("01", "Chat — your PE AI squad on call",
         "triage, LBO modelling, IC memo, VDR audit"),
        ("02", "Pipeline — kanban across every deal stage",
         "sector + ownership filters, click into a deal"),
        ("03", "Deal detail — brief on the right, chat in the centre",
         "LTM, top customers, DD findings, LBO returns"),
        ("04", "Analytics — ask in English, get a chart",
         "sector multiples, stage counts, LP mix"),
        ("05", "Instructions — tune how each specialist thinks",
         "edit in-app, changes land on the next conversation"),
        ("06", "Extensions — live web + local registries",
         "sourcing beyond your CRM, Baltic registry lookups"),
    ]

    y = Cm(5.4)
    row_h = Cm(1.8)
    for idx, (num, title, sub) in enumerate(rows):
        _text(s, Cm(1.5), y + idx * row_h, Cm(1.8), row_h,
              num, size=18, bold=True, color=ACCENT)
        _text(s, Cm(4), y + idx * row_h, Cm(14), row_h,
              title, size=14, bold=True, color=INK)
        _text(s, Cm(18.5), y + idx * row_h, SLIDE_W - Cm(20), row_h,
              sub, size=12, color=INK_MUTED)

    return s


def _content_slide(prs, *, eyebrow, title, subtitle, bullets, screenshot, caption=None):
    s = _blank(prs)
    # Left column: text
    left_x = Cm(1.5)
    left_w = Cm(13)

    _text(s, left_x, Cm(1.5), left_w, Cm(0.8),
          eyebrow.upper(), size=10, bold=True, color=ACCENT)
    _text(s, left_x, Cm(2.2), left_w, Cm(2.2),
          title, size=22, bold=True, color=INK)
    _text(s, left_x, Cm(4.3), left_w, Cm(1.3),
          subtitle, size=13, color=INK_MUTED)
    _bullets(s, left_x, Cm(6.2), left_w, Cm(10),
             bullets, size=12)

    # Right column: screenshot
    shot_path = SHOTS / screenshot
    if shot_path.exists():
        max_w_cm = 16.0
        max_h_cm = 13.0
        w_in, h_in = _fit_image_inches(shot_path,
                                       max_w_in=max_w_cm / 2.54,
                                       max_h_in=max_h_cm / 2.54)
        img_x = SLIDE_W - Cm(1.5) - Inches(w_in)
        img_y = Cm(2.5)
        s.shapes.add_picture(str(shot_path), img_x, img_y,
                             width=Inches(w_in), height=Inches(h_in))
        if caption:
            _text(s, img_x, img_y + Inches(h_in) + Cm(0.2),
                  Inches(w_in), Cm(0.8),
                  caption, size=9, color=INK_DIM, align=PP_ALIGN.CENTER,
                  font="Inter")
    else:
        _text(s, Cm(16), Cm(8), SLIDE_W - Cm(17.5), Cm(2),
              f"[missing {screenshot}]", size=11, color=INK_DIM)

    return s


def _closing_slide(prs):
    s = _blank(prs)
    _text(s, Cm(0), Cm(4.5), SLIDE_W, Cm(1),
          "LET'S TALK", size=11, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
    _text(s, Cm(0), Cm(5.8), SLIDE_W, Cm(3),
          "See PEHero on your deals.",
          size=40, bold=True, color=INK, align=PP_ALIGN.CENTER)
    _text(s, Cm(3), Cm(9.5), SLIDE_W - Cm(6), Cm(3),
          "Book a 20-minute walkthrough. We'll load one of your recent deals into "
          "PEHero and show you the full agent flow — live.",
          size=16, color=INK_MUTED, align=PP_ALIGN.CENTER)
    _text(s, Cm(0), Cm(13.5), SLIDE_W, Cm(1),
          "hello@pehero.fyi   ·   pehero.fyi/contact",
          size=14, bold=True, color=INK, align=PP_ALIGN.CENTER)
    _text(s, Cm(0), Cm(15), SLIDE_W, Cm(1),
          "BYOD — bring your own deal data.",
          size=10, color=INK_DIM, align=PP_ALIGN.CENTER)
    return s


def build() -> None:
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    _hero_slide(prs)
    _agenda_slide(prs)

    # Chat section
    _content_slide(
        prs,
        eyebrow="01 · Chat",
        title="One chat, every PE specialist",
        subtitle="Type a prefix or plain English — PEHero picks the right agent.",
        bullets=[
            "Left: your sessions, the full agent squad, and Pipeline / Instructions / Analytics.",
            "Centre: the conversation. Contextual sample prompts appear right under the input.",
            "Right: tables, citations, memo previews stream in as the agent works.",
            "A live 'thinking' indicator shows what's happening behind the scenes.",
        ],
        screenshot="07-chat-empty.png",
        caption="Empty chat with contextual prompts",
    )
    _content_slide(
        prs,
        eyebrow="01 · Chat",
        title="Deal Triage — go / no-go in 90 seconds",
        subtitle="A quick decision backed by comps and market signals.",
        bullets=[
            "Type 'triage:' or describe the deal in plain English.",
            "The agent pulls comparables and sector context on its own.",
            "Returns a verdict, a three-bullet rationale, and a concrete next step.",
            "A follow-up button surfaces to take the next step in one click.",
        ],
        screenshot="08-chat-triage.png",
        caption="Live triage on a vertical-SaaS target",
    )
    _content_slide(
        prs,
        eyebrow="01 · Chat",
        title="LBO Model Builder",
        subtitle="Full 5-year model and sensitivity from one sentence of assumptions.",
        bullets=[
            "Normalises seller financials with standard QoE add-backs.",
            "Projects revenue, margin, capex, interest, free cash flow and debt paydown.",
            "Returns IRR, MOIC and a value-creation bridge — kept on hand for re-use.",
            "Year-by-year table appears in the right pane, editable by asking for changes.",
        ],
        screenshot="09-chat-lbo.png",
        caption="LBO built from a plain-English prompt",
    )
    _content_slide(
        prs,
        eyebrow="01 · Chat",
        title="IC Memo Writer",
        subtitle="An investment-committee memo, drafted from the deal's own data.",
        bullets=[
            "Pulls the deal brief, LTM financials, LBO model, debt stack, comps and findings.",
            "Drafts thesis, market, financials, value-creation plan, risks, recommendation.",
            "Every quantitative claim is sourced from the deal — no invented numbers.",
            "IC length by default; ask for a one-pager and it re-writes accordingly.",
        ],
        screenshot="10-chat-memo.png",
        caption="IC memo generated from the deal data",
    )

    # Pipeline section
    _content_slide(
        prs,
        eyebrow="02 · Pipeline",
        title="Kanban across every deal stage",
        subtitle="Sourced → Closed / Held / Exited — every live target on one board.",
        bullets=[
            "Each card shows sector, LTM revenue and EBITDA, ask EV and multiple.",
            "A heat dot on the card reflects seller intent — cold, warm or hot.",
            "Sector and ownership chips filter the board in a click.",
            "Click a card to open the full deal workspace.",
        ],
        screenshot="11-pipeline-kanban.png",
        caption="Full pipeline kanban",
    )
    _content_slide(
        prs,
        eyebrow="02 · Pipeline",
        title="Filter to what matters",
        subtitle="Narrow to a sector or ownership type in one click.",
        bullets=[
            "Filter chips: sector, and ownership (founder, family, PE, VC, carve-out).",
            "Cards update instantly without a page jump.",
            "Perfect for mandate conversations — 'show me lower-mid-market software only'.",
            "Filtered views are shareable with a single URL.",
        ],
        screenshot="12-pipeline-software.png",
        caption="Pipeline filtered to software",
    )
    _content_slide(
        prs,
        eyebrow="03 · Deal detail",
        title="Every deal has its own workspace",
        subtitle="Brief on the right, chat in the centre, artifacts stream in as the squad works.",
        bullets=[
            "Right pane: HQ, LTM financials, top customers, DD findings, margin and multiple.",
            "Centre: per-deal chat — 'triage this', 'draft the IC memo', 'show DD findings'.",
            "Any specialist can be invoked without ever leaving the deal.",
            "New artifacts from tool calls land alongside the brief as they arrive.",
        ],
        screenshot="13-pipeline-deal.png",
        caption="Single-deal workspace",
    )

    # Analytics
    _content_slide(
        prs,
        eyebrow="04 · Analytics",
        title="Ask in English, get a chart",
        subtitle="Analytics that read the same data your deal team does.",
        bullets=[
            "Natural-language questions run read-only against your deal data.",
            "The right chart and title are picked automatically from the result.",
            "Curated sample questions seed the experience for first-time users.",
            "The underlying query is shown under every chart — fully auditable.",
        ],
        screenshot="15-analytics-stages.png",
        caption="Company count by deal stage",
    )
    _content_slide(
        prs,
        eyebrow="04 · Analytics",
        title="Sector multiples over time",
        subtitle="Median EV/EBITDA by sector, rolling 24 months.",
        bullets=[
            "Market signals cover every sector and sub-sector you're tracking.",
            "A grouped line chart, colour-coded by sector, appears in one click.",
            "Answers 'what's happening to multiples in X' without a spreadsheet.",
            "Drops straight into the LP update or IC pre-read.",
        ],
        screenshot="16-analytics-sector.png",
        caption="Median EV/EBITDA by sector",
    )

    # Instructions
    _content_slide(
        prs,
        eyebrow="05 · Instructions",
        title="Tune the squad, live",
        subtitle="Every specialist's instructions are editable — from the same interface.",
        bullets=[
            "Each role has its own set of instructions, plus a shared PE glossary.",
            "Edits save in-place and take effect on the very next conversation.",
            "No restarts, no deploys. Just change how the squad thinks and carry on.",
            "Perfect for onboarding a partner's preferred memo style or diligence approach.",
        ],
        screenshot="17-instructions-list.png",
        caption="The full squad, editable",
    )
    _content_slide(
        prs,
        eyebrow="05 · Instructions",
        title="Editing the Deal Triage instructions",
        subtitle="Workflow, tone and output format — all in one view.",
        bullets=[
            "Rewrite a specialist's instructions the way you'd brief a new associate.",
            "Shared PE context is applied automatically so you don't repeat yourself.",
            "Versioned alongside the product so changes are auditable.",
            "Great for encoding your house style once and letting it apply everywhere.",
        ],
        screenshot="18-instructions-edit.png",
        caption="Editing a specialist's instructions",
    )

    _closing_slide(prs)

    OUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(OUT)
    size_kb = OUT.stat().st_size / 1024
    print(f"Wrote {OUT}  ({size_kb:.1f} KB)")


if __name__ == "__main__":
    build()
